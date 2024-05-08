import pandas as pd # type: ignore
from datetime import datetime
import customtkinter as ctk
from customtkinter import filedialog
from tkinter import messagebox
import winreg
import zipfile
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from os import listdir,walk,makedirs
from os.path import join,relpath,isdir,exists, isfile
import img2pdf
import tempfile

TYPE=['.pdf', 'PPT','.cbz']

F_type=['Folder', 'Folder > Folder', 'Single', 'PPT']

def get_windows_downloads_folder():
    try:
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, r"Software\\Microsoft\Windows\\CurrentVersion\\Explorer\Shell Folders") as key:
            downloads_guid = '{374DE290-123F-4565-9164-39C4925E467B}' # every single built it folder has a GUID
            downloads_folder, _ = winreg.QueryValueEx(key, downloads_guid)
            return downloads_folder
    except Exception as e:
        print("Error:", e)
        return None

class App(ctk.CTk):
    global r
    r=0
    class MyScrollableCheckboxFrame(ctk.CTkScrollableFrame):
        def __init__(self, master, title, values):
            super().__init__(master, label_text=title)
            self.grid_columnconfigure(0, weight=5)
            self.grid_columnconfigure(1, weight=1)
            self.grid_columnconfigure(2, weight=1)
            self.grid_columnconfigure(3, weight=1)
            self.values = values
            self.rows = []

            for row,data in enumerate(values):
                combobox = ctk.CTkOptionMenu(self,values=TYPE)
                combobox2 = ctk.CTkOptionMenu(self,values=F_type)
                # button=ctk.CTkButton(self,text="X",command= lambda: self.removedata())
                label = ctk.CTkLabel(self, text=data)
                label.grid(row=row, column=0, padx=10, pady=(10, 0), sticky="w")
                combobox.grid(row=row, column=1, padx=10, pady=(10, 0), sticky="w")
                combobox2.grid(row=row, column=2, padx=10, pady=(10, 0), sticky="w")
                # button.grid(row=row, column=3, padx=10, pady=(10, 0), sticky="w")
                self.rows.append({'label':label,'combo':combobox,'folder':combobox2})

        def get(self):
            all_data = []
            for row in self.rows:
                all_data.append({'dst':row['label'].cget("text"),"type":row['combo'].get(),'f_type':row['folder'].get()})
            return all_data
        
        def removedata(self):
            # self.rows.pop(index)
            # TODO no idea
            # self.grid_remove(row=index)
            # self.add(combobox)
            pass

    def __init__(self):
        super().__init__()
        self.title("Image to PDF")
        self.geometry("600x600")
        self.grid_columnconfigure(0, weight=1)
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=5)
        self.grid_rowconfigure(2, weight=1)
        self.grid_rowconfigure(3, weight=1)

        self.button=ctk.CTkButton(self,text="Add Folder", command=self.button_callbck)
        self.button.grid(row=0, column=0, sticky="w",padx=10, pady=(10, 0))
        self.values = []
        self.scrollable_checkbox_frame = self.MyScrollableCheckboxFrame(self, title="Folder list", values=self.values)
        self.scrollable_checkbox_frame.grid(row=1, column=0, padx=10, pady=(10, 0), sticky="nsew",columnspan=2)
        self.outputlfolder=ctk.CTkLabel(self, text=get_windows_downloads_folder())
        self.outputlfolder.grid(row=2, column=0, sticky="ns",padx=10, pady=(10, 0))
        self.button3=ctk.CTkButton(self,text="Output Folder",command=self.set_output)
        self.button3.grid(row=2, column=1, sticky="w",padx=10, pady=(10, 0))
        self.button2=ctk.CTkButton(self,text="Convert", command=self.convert)
        self.button2.grid(row=3, column=0, sticky="w",padx=10, pady=(10, 0))

    def button_callbck(self):
        folder_selected = filedialog.askdirectory()
        if len(folder_selected)!=0:
            # self.scrollable_checkbox_frame.destroy()
            self.values.append(folder_selected)
            self.scrollable_checkbox_frame = self.MyScrollableCheckboxFrame(self, title="Folder list", values=self.values)
            self.scrollable_checkbox_frame.grid(row=1, column=0, padx=10, pady=(10, 0), sticky="nsew",columnspan=2)

    def set_output(self):
        folder_selected = filedialog.askdirectory()
        self.outputlfolder.configure(text=folder_selected)

    def create_cbz(self, source_folder, output_cbz):
        start_time = datetime.now()
        with zipfile.ZipFile(output_cbz, 'w') as cbz_file:
            num_files = 0
            for root, _, files in walk(source_folder):
                for file in files:
                    if file.lower().endswith(('.jpg', '.jpeg', '.png', '.gif')):
                        num_files += 1
                        image_path = join(root, file)
                        if not isdir(image_path):
                            cbz_file.write(image_path, relpath(image_path, source_folder))
        conversion_time = datetime.now() - start_time

        print('output_cbz: ', output_cbz)

        # Log conversion information
        self.log_conversion_info(source_folder, output_cbz, '.cbz', num_files, conversion_time.total_seconds())

    def create_pdf(self,source_folder, output_pdf):
        imgs=[]
        a4inpt = (img2pdf.mm_to_pt(210),img2pdf.mm_to_pt(297))
        layout_fun = img2pdf.get_layout_fun(a4inpt)
        for fname in listdir(source_folder):
            if fname.endswith(('.jpg', '.jpeg', '.png', '.gif')):
                path=join(source_folder,fname)
                if not isdir(path):
                    imgs.append(path)
            # elif fname.endswith(('.cbz')):
            #     self.cbztopdf(source_folder,output_pdf)

        start_time = datetime.now()
        with open(output_pdf,"wb") as f:#take folder name as pdf name
            f.write(img2pdf.convert(imgs,layout_fun=layout_fun))
        conversion_time = datetime.now() - start_time
        num_pages = len(imgs)

        print('\033[91m'+'output_pdf: ' + '\033[92m', output_pdf)
        # ctk.CTkFrame(master=app, text='Conversion Completed')
        messagebox.showinfo("Conversion Completed", "The conversion process has been successfully completed.")

        imgs.clear()

        self.log_conversion_info(source_folder, output_pdf, '.pdf', num_pages, conversion_time.total_seconds())

    def create_ppt_from_folder(self, folder_path, output_ppt_path):
        prs = Presentation()
        num_slides = 0
        start_time = datetime.now()
        for file in listdir(folder_path):
            file_path = join(folder_path, file)
            if isfile(file_path) and file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
                img = Image.open(file_path)
                width_px, height_px = img.size
                img.close()

                # Convert pixels to inches (assuming 96 dpi)
                width_in = width_px / 96
                height_in = height_px / 96

                # Calculate the center position
                slide_layout = prs.slide_layouts[5]  # Using a blank slide layout
                slide = prs.slides.add_slide(slide_layout)
                left_in = (prs.slide_width - Inches(width_in)) / 2
                top_in = (prs.slide_height - Inches(height_in)) / 2

                # Add image to slide and center it
                slide.shapes.add_picture(file_path, left_in, top_in, Inches(width_in), Inches(height_in))
                num_slides += 1
        prs.save(output_ppt_path)
        conversion_time = datetime.now() - start_time
        self.log_conversion_info(folder_path, output_ppt_path, 'PPT', num_slides, conversion_time.total_seconds())

    def log_conversion_info(self, source_folder, output_file, c_type, num_pages, conversion_time):
        current_date = datetime.now().strftime("%d/%m/%y")
        data = {
            'Date': [current_date],
            'Source Folder': [source_folder],
            'Output File': [output_file],
            'Type': [c_type],
            'Number of Pages': [num_pages],
            'Conversion Time (seconds)': [conversion_time]
        }
        df = pd.DataFrame(data)
        df.to_csv('conversion_log.csv', mode='a', header=not exists('conversion_log.csv'), index=False)


    def cbztopdf(self,source_folder, output_pdf):
        temp_dir = tempfile.TemporaryDirectory()
        print('\033[91m'+'temp_dir.name: ' + '\033[92m', temp_dir.name)
        with zipfile.ZipFile(source_folder, 'r') as zip_ref:
            zip_ref.extractall(temp_dir.name)
            l=listdir(temp_dir)
            print('\033[91m'+'l: ' + '\033[92m', l)

    def folderInsideFolder(self,source_folder:str,output:str,c_type:str):
        for i in listdir(source_folder):
            source_folder=join(source_folder,i)
            if isdir(source_folder):
                self.onlyFolder(source_folder,output,i,c_type)
            else:
                print('\033[91m'+'single_file: ' + '\033[92m', source_folder)

    def onlyFolder(self,source_folder:str,output:str,filename:str,c_type:str):
        output=self.makeFile(c_type,filename,output)
        print('\033[91m'+'output: ' + '\033[92m', output)
        if c_type==".cbz":
            self.create_cbz(source_folder,output)
        elif c_type==".pdf":
            self.create_pdf(source_folder,output)

    def SingleFile(self,source_folder:str,output:str,filename:str,c_type:str):
        output=self.makeFile(c_type,filename,output)
        print('\033[91m'+'output: ' + '\033[92m', output)
        if c_type==".cbz":
            self.create_cbz(source_folder,output)
        elif c_type==".pdf":
            self.create_pdf(source_folder,output)

    def makeFile(self,filetype,filename,output):
        full_path=join(output,filename+filetype)
        while(True):
            if not exists(full_path):
                return full_path
            else:
                full_path=join(output,filename+' -copy'+filetype)
                print('\033[91m'+'full_path: ' + '\033[92m', full_path)

    def makeFolder(self,folder_name,output):
        full_path=join(output,folder_name)
        while(True):
            if not exists(full_path):
                makedirs(full_path)
                return full_path
            else:
                full_path+=' -copy'
                print('\033[91m'+'full_path: ' + '\033[92m', full_path)

    def convert(self):
       data=self.scrollable_checkbox_frame.get()
       print('\033[91m'+'data: ' + '\033[92m', data)
       output=self.outputlfolder.cget("text")
       for i in data:
            dst=i['dst']
            c_type=i['type']
            f_type=i['f_type']
            title=dst.split('/')[-1].split('.')[0]
            if f_type=='Folder > Folder':
                output=self.makeFolder(title,output)
                self.folderInsideFolder(dst,output,c_type)
            elif f_type=='Folder':
                self.onlyFolder(dst,output,title,c_type)
            elif f_type == 'PPT':  # New option for PPT conversion
                output_ppt_path = join(output, title + '.pptx')
                self.create_ppt_from_folder(dst, output_ppt_path)
                messagebox.showinfo("Conversion Completed", "The conversion process has been successfully completed.")


if __name__=='__main__':
    app = App()
    app.mainloop()